from pathlib import Path

from simon.config.settings import settings
from simon.knowledge.embeddings import OpenAIEmbeddings
from simon.knowledge.retrieval import FileRetriever


class KnowledgeBase:
    """Document ingestion + retrieval with hidden implementation details."""

    def __init__(
        self,
        chunk_size: int = 500,
        overlap: int = 50,
        store_path: str | None = None,
    ) -> None:
        path = store_path or settings.knowledge_store_path or ".simon_knowledge"
        self._retriever = FileRetriever(OpenAIEmbeddings(), path)
        self.chunk_size = chunk_size
        self.overlap = overlap

    def add(self, path: str) -> int:
        p = Path(path)
        if not p.exists():
            raise FileNotFoundError(f"Knowledge path not found: {path}")

        files = [p] if p.is_file() else [f for f in p.rglob("*") if f.is_file()]
        count = 0
        for file in files:
            text = self._read_file(file)
            chunks = self._chunk(text)
            if chunks:
                self._retriever.add_source(chunks, source=str(file))
                count += len(chunks)
        return count

    def search(self, query: str, top_k: int = 3) -> list[dict[str, str]]:
        return self._retriever.search(query, top_k=top_k)

    def _chunk(self, text: str) -> list[str]:
        if not text.strip():
            return []

        chunks: list[str] = []
        step = max(1, self.chunk_size - self.overlap)
        for i in range(0, len(text), step):
            chunk = text[i : i + self.chunk_size].strip()
            if chunk:
                chunks.append(chunk)
        return chunks

    def _read_file(self, path: Path) -> str:
        ext = path.suffix.lower()

        if ext == ".pdf":
            try:
                from pypdf import PdfReader
            except ImportError as exc:
                raise RuntimeError("Install pypdf to ingest PDF files.") from exc
            reader = PdfReader(str(path))
            return "\n".join(page.extract_text() or "" for page in reader.pages)

        if ext == ".docx":
            try:
                import docx
            except ImportError as exc:
                raise RuntimeError("Install python-docx to ingest Word files.") from exc
            doc = docx.Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs)

        if ext == ".xlsx":
            try:
                import openpyxl
            except ImportError as exc:
                raise RuntimeError("Install openpyxl to ingest Excel files.") from exc
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            lines = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    line = "\t".join("" if v is None else str(v) for v in row)
                    if line.strip():
                        lines.append(line)
            return "\n".join(lines)

        if ext == ".pptx":
            try:
                from pptx import Presentation
            except ImportError as exc:
                raise RuntimeError("Install python-pptx to ingest PowerPoint files.") from exc
            prs = Presentation(str(path))
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        lines.append(shape.text_frame.text)
            return "\n".join(lines)

        return path.read_text(encoding="utf-8", errors="ignore")
